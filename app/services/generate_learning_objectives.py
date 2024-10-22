import os, sys, csv
from openAI_client_base import OpenAIClientBase
from openai import OpenAI
from pptx import Presentation
from pathlib import Path

class Generate_Learning_Objectives:

    def __init__(self):
        super().__init__()

    def get_files(self, input_path='IHL Hypersensitivities Fall 2024.pptx'):
        path = Path(input_path)
        output_prefix = path.stem
        output_file = os.path.join(os.getcwd(), f'{output_prefix} Learning Objectives.csv') # TODO: don't love formatting....
        if path.is_file():
            pptx_files = [input_path]
        elif path.is_dir():
            pptx_files = list(path.glob('*.pptx'))
        else:
            print("The provided path is not a valid file or directory.")
            sys.exit(1)
        return pptx_files, output_file

    def extract_text_from_pptx(self, pptx_file):
        prs = Presentation(pptx_file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    def define_objectives_from_pptx(self, pptx_file, temperature=1.0):
        text = self.extract_text_from_pptx(pptx_file)
        print(text)


    def main(self):
        pptx_files, output_file = self.get_files()
        for pptx_file in pptx_files:
            objectives = self.define_objectives_from_pptx(pptx_file)
        # with open(output_file, 'a', newline='', encoding='utf-8') as csvfile:
        #     csv_writer = csv.writer(csvfile)
        #     csv_writer.writerow(['name', 'learning_objective','tokens','emb'])

        #     for pptx_file in pptx_files:
        #         objectives = self.define_objectives_from_pptx(pptx_file)
                # tag = Path(pdf_file).stem
                # write_to_csv(csv_writer, tag, objectives)


learning_objectives_generator = Generate_Learning_Objectives()
learning_objectives_generator.main()